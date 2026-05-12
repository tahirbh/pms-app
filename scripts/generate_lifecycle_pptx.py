import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_pms_lifecycle_pptx():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Global PMS: Complete Application Lifecycle"
    subtitle.text = "Building a Robust Property Management Ecosystem\nVersion 1.0.39"

    # Slide 2: The Core Lifecycle (Infographic)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Ecosystem Lifecycle Infographic"

    # Define steps
    steps = [
        {"name": "1. ASSET SETUP", "desc": "Property Registration & Annual Rent Mapping"},
        {"name": "2. CONTRACTING", "desc": "Hijri/Gregorian Tenancy & Payment Plans"},
        {"name": "3. LEDGER GEN", "desc": "Automated Installment Creation & Pro-rata"},
        {"name": "4. OPERATIONS", "desc": "Income/Expense Tracking & Maintenance"},
        {"name": "5. ANALYTICS", "desc": "Dashboard Drill-down & Pivot Reporting"}
    ]

    left = Inches(0.5)
    top = Inches(2)
    width = Inches(1.5)
    height = Inches(1.5)

    for i, step in enumerate(steps):
        # Create Circle
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        shape.fill.solid()
        shape.fill.foreground_color.rgb = RGBColor(37, 99, 235) # Primary Blue
        
        # Add text to circle
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = step["name"]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Add description below
        txBox = slide.shapes.add_textbox(left, top + height + Inches(0.2), width, Inches(1))
        tf = txBox.text_frame
        p2 = tf.paragraphs[0]
        p2.text = step["desc"]
        p2.font.size = Pt(10)
        
        left += Inches(1.8)

    # Slide 3: Technical Scope Matrix
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Technical Scope & Capabilities"

    rows = 6
    cols = 2
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(4)).table
    
    table.cell(0, 0).text = "Domain"
    table.cell(0, 1).text = "Implementation Detail"
    
    data = [
        ("Calendar", "Dynamic Hijri/Gregorian pro-rata logic"),
        ("Data Layer", "Supabase Real-time DB with RLS Security"),
        ("Financials", "Ledger-based income/expense tracking"),
        ("Support", "Admin Support Mode & Team Impersonation"),
        ("UX", "Glassmorphism responsive Design System")
    ]
    
    for i, (domain, detail) in enumerate(data):
        table.cell(i+1, 0).text = domain
        table.cell(i+1, 1).text = detail

    # Save
    file_path = os.path.join(os.getcwd(), "Global_PMS_Lifecycle_Scope.pptx")
    prs.save(file_path)
    print(f"PPTX generated at: {file_path}")

if __name__ == "__main__":
    create_pms_lifecycle_pptx()
